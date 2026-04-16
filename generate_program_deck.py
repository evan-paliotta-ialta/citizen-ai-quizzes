#!/usr/bin/env python3
"""
Generate: Citizen AI Engineer Program — Master Deck.pptx
19 slides covering the full program: policy, tooling, compliance, workflows,
Highlander tracking, course, cadence, and costs.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ─── Paths ────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(os.path.dirname(BASE), "Pitch", "ialta-logo-blue.png")
OUT  = os.path.join(BASE, "Citizen AI Engineer Program — Master Deck.pptx")

# ─── Brand colours ────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x11, 0x17)   # near-black background
BLUE      = RGBColor(0x2D, 0x6F, 0xF7)   # iAltA blue
BLUE_DK   = RGBColor(0x1A, 0x48, 0xBB)   # darker blue (table header alt)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xCC, 0xD1, 0xD9)   # light-grey body text
MGRAY     = RGBColor(0x55, 0x60, 0x6E)   # mid-grey muted
ROW_ALT   = RGBColor(0x1A, 0x1F, 0x2E)   # dark-blue alternating row
ROW_EVEN  = RGBColor(0x13, 0x16, 0x20)   # slightly darker even row
ACCENT_R  = RGBColor(0xE8, 0x4A, 0x3A)   # red for no-drive zone
ACCENT_G  = RGBColor(0x27, 0xAE, 0x60)   # green for good-faith col
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # gold for future state

# ─── Slide dimensions (16:9) ──────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    # fill background
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def box(sl, x, y, w, h):
    return sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def rect(sl, x, y, w, h, fill_rgb, line_rgb=None):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
    else:
        shp.line.fill.background()
    return shp


def add_run(para, text, size=18, bold=False, color=None, italic=False):
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    run.font.name  = "Calibri"
    return run


def para_in(tf, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after  = Pt(0)
    return p


def label(sl, text, x, y, w, h=0.35, size=9, color=None, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = box(sl, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    add_run(p, text, size=size, color=color or LGRAY, bold=bold, italic=italic)
    return tb


def section_pill(sl, text, x=0.3, y=0.18):
    """Small coloured pill above a slide title."""
    r = rect(sl, x, y, len(text) * 0.095 + 0.3, 0.28, BLUE)
    tf = r.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text.upper(), size=8, bold=True, color=WHITE)


def slide_title(sl, text, y=0.55, size=30, x=0.3, w=12.7):
    tb = box(sl, x, y, w, 0.7)
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    add_run(p, text, size=size, bold=True, color=WHITE)


def divider_line(sl, y=1.25, x=0.3, w=12.73):
    ln = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.025))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()


def bullet_block(sl, items, x, y, w, h, size=13, gap=0.32, title=None, title_size=14):
    tb = box(sl, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else para_in(tf)
        first = False
        add_run(p, title, size=title_size, bold=True, color=BLUE)
        para_in(tf)  # spacer
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else para_in(tf)
        first = False
        p.space_after = Pt(4)
        if item.startswith("**") and "**" in item[2:]:
            # bold prefix pattern: **bold** rest
            parts = item.split("**")
            add_run(p, "• ", size=size, color=BLUE, bold=True)
            for j, part in enumerate(parts):
                if part:
                    add_run(p, part, size=size, bold=(j % 2 == 1), color=WHITE if j % 2 == 0 else WHITE)
        else:
            add_run(p, "• " + item, size=size, color=LGRAY)


# ─── Table builder ────────────────────────────────────────────────────────────

def make_table(sl, rows, cols, x, y, w, h, col_widths=None):
    tbl = sl.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    return tbl


def cell_fill(cell, rgb):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def cell_text(cell, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.06)
    tf.margin_right  = Inches(0.06)
    tf.margin_top    = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    add_run(p, text, size=size, bold=bold, color=color or LGRAY, italic=italic)


def header_row(tbl, headers, size=11, bg=BLUE_DK):
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        cell_fill(c, bg)
        cell_text(c, h, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def data_rows(tbl, data, start_row=1, size=10, alt=True):
    for ri, row_data in enumerate(data):
        bg = ROW_ALT if (ri % 2 == 0 and alt) else ROW_EVEN
        for ci, val in enumerate(row_data):
            c = tbl.cell(start_row + ri, ci)
            cell_fill(c, bg)
            cell_text(c, val, size=size)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()

# Full-width blue bar at top
rect(sl, 0, 0, 13.33, 0.06, BLUE)

# Logo (if it exists)
if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, Inches(0.35), Inches(0.25), width=Inches(1.6))

# Main title
tb = box(sl, 1.5, 2.6, 10.3, 1.1)
tf = tb.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Citizen AI Engineer Program", size=44, bold=True, color=WHITE)

# Subtitle
tb2 = box(sl, 1.5, 3.75, 10.3, 0.5)
tf2 = tb2.text_frame
p2  = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
add_run(p2, "From Policy to Practice — A Complete Operating Model", size=18, color=LGRAY, italic=True)

# Bottom bar
rect(sl, 0, 7.2, 13.33, 0.3, BLUE)
tb3 = box(sl, 0, 7.18, 13.33, 0.3)
tf3 = tb3.text_frame
p3  = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
add_run(p3, "iAltA  |  Confidential  |  2026", size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT & WHY
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Program Overview")
slide_title(sl, "What & Why")
divider_line(sl)

# Four pillars
pillars = [
    ("Policy", "Acceptable Use Policy + Driving Zones framework govern what's allowed and when"),
    ("Tooling", "Claude Teams (now) + GitHub Teams — enterprise AI + version-controlled audit trail"),
    ("Tracking", "Highlander ODL ingests every commit — proves business value to CIO and board"),
    ("Culture", "Driver's license model: certification + signed AUP before any license is issued"),
]
pw = 2.9
px = 0.35
for i, (title, body) in enumerate(pillars):
    bx = px + i * (pw + 0.18)
    rect(sl, bx, 1.55, pw, 0.06, BLUE)
    r = rect(sl, bx, 1.61, pw, 4.1, ROW_ALT)
    tb = box(sl, bx + 0.12, 1.75, pw - 0.24, 3.7)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, title, size=16, bold=True, color=BLUE)
    p2 = para_in(tf)
    p2.space_before = Pt(6)
    add_run(p2, body, size=12, color=LGRAY)

# ROI callout
r2 = rect(sl, 0.35, 5.9, 12.63, 0.9, RGBColor(0x0A, 0x25, 0x5A))
tb4 = box(sl, 0.45, 5.95, 12.43, 0.8)
tf4 = tb4.text_frame
p4  = tf4.paragraphs[0]
p4.alignment = PP_ALIGN.CENTER
add_run(p4, "25 citizens × 1 hr/week saved × $75/hr = ", size=15, color=LGRAY)
add_run(p4, "$93,750/year in recaptured capacity", size=15, bold=True, color=WHITE)
p4b = para_in(tf4)
p4b.alignment = PP_ALIGN.CENTER
add_run(p4b, "Highlander measures and proves every dollar of it", size=11, color=BLUE, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROGRAM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Program Overview")
slide_title(sl, "Program Architecture — How It All Connects")
divider_line(sl)

steps = [
    ("AUP &\nPolicy", BLUE),
    ("License\nIssued", RGBColor(0x1A, 0x48, 0xBB)),
    ("Certif-\nication", BLUE),
    ("Claude\nTooling", RGBColor(0x1A, 0x48, 0xBB)),
    ("GitHub\nTeams", BLUE),
    ("Highlander\nODL", RGBColor(0x1A, 0x48, 0xBB)),
    ("OKR\nProof", BLUE),
]
bw, bh, by = 1.55, 1.0, 2.4
total_w = len(steps) * bw + (len(steps) - 1) * 0.25
start_x = (13.33 - total_w) / 2

for i, (txt, clr) in enumerate(steps):
    bx = start_x + i * (bw + 0.25)
    rect(sl, bx, by, bw, bh, clr)
    tb = box(sl, bx, by, bw, bh)
    tf = tb.text_frame
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, txt, size=13, bold=True, color=WHITE)
    # Arrow
    if i < len(steps) - 1:
        ax = bx + bw + 0.02
        label(sl, "→", ax, by + 0.35, 0.22, size=18, color=BLUE, bold=True)

# Subtext under each box
subtexts = [
    "Governs all use",
    "Course + signature",
    "16-module cert",
    "Desktop or Code",
    "Audit trail",
    "Metrics engine",
    "Board narrative",
]
for i, sub in enumerate(subtexts):
    bx = start_x + i * (bw + 0.25)
    label(sl, sub, bx, by + bh + 0.08, bw, size=9, color=MGRAY, align=PP_ALIGN.CENTER)

# Zone tiers callout
label(sl, "Zone 1 → Zone 2 → Zone 3 → No-Drive   |   The data type determines which zone applies — not the tool used", 0.35, 5.7, 12.63, size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
label(sl, "Every citizen commit flows into Highlander automatically. No manual reporting required.", 0.35, 6.1, 12.63, size=11, color=LGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI ACCEPTABLE USE POLICY
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Governance & Policy")
slide_title(sl, "AI Acceptable Use Policy")
divider_line(sl)

label(sl, "The document that makes us defensible.  Required by SOC2 CC6.1 (access control) and ISO 27001 A.5.10 (acceptable use).", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

aup_items = [
    ("ALLOWED — Zone 1", "Public data only. No company information. No approval needed. Just use Claude.", WHITE, ROW_ALT),
    ("ALLOWED WITH PROCESS — Zone 2", "Internal non-confidential data. GitHub repo required. Code review before any merge.", WHITE, RGBColor(0x0D, 0x2A, 0x55)),
    ("ALLOWED WITH FULL PROCESS — Zone 3", "Confidential / PII / client data. Requires: Jira ticket + approved MCP only + manager sign-off + full audit trail.", WHITE, RGBColor(0x0A, 0x1F, 0x44)),
    ("PROHIBITED — No-Drive Zone", "Credentials, secrets, bypassing guardrails, unapproved tools, consumer claude.ai with company data.", ACCENT_R, RGBColor(0x2A, 0x08, 0x08)),
    ("PROHIBITED — Unapproved AI tools", "Consumer ChatGPT, Gemini, or any non-enterprise AI tool with company data. Claude Teams is the only approved AI.", ACCENT_R, RGBColor(0x2A, 0x08, 0x08)),
    ("VIOLATION CONSEQUENCE", "First offense: mandatory re-certification. Second offense: license revocation. Intentional bypass: disciplinary action.", GOLD, RGBColor(0x2A, 0x1A, 0x00)),
]

ty = 1.75
for zone, body, text_clr, bg_clr in aup_items:
    rect(sl, 0.3, ty, 12.73, 0.62, bg_clr)
    label(sl, zone, 0.42, ty + 0.06, 3.1, size=10, bold=True, color=text_clr)
    label(sl, body, 3.6, ty + 0.06, 9.3, size=10, color=LGRAY)
    ty += 0.66

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DRIVING ZONES MODEL
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Governance & Policy")
slide_title(sl, "The Driving Zones Model")
divider_line(sl)

label(sl, "The data determines the zone. The zone determines the process. No gray areas.", 0.3, 1.35, 12.73, size=13, color=BLUE, bold=True)

tbl = make_table(sl, 5, 4, 0.3, 1.72, 12.73, 4.6, col_widths=[2.2, 2.1, 3.2, 5.23])
header_row(tbl, ["Zone", "Analogy", "Data Involved", "Process Required"])

zones = [
    ("Zone 1\nPrivate Property", "Your driveway", "Public data only. No company information.", "None — just use Claude."),
    ("Zone 2\nResidential", "Neighborhood streets", "Internal, non-confidential. No PII or client data.", "GitHub repo + branch protection + code review before merge."),
    ("Zone 3\nHighway", "Highway driving", "Confidential / PII / client data, MCP servers, production systems.", "Zone 2 + Jira ticket + approved MCP servers only + audit logging + manager review."),
    ("No-Drive Zone\nZero Tolerance", "Drunk driving", "Credentials, secrets, bypassing guardrails, unapproved tools.", "Prohibited always. Violation = license revocation."),
]
for ri, row in enumerate(zones):
    bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        if ri == 3:
            cell_fill(c, RGBColor(0x2A, 0x08, 0x08))
            cell_text(c, val, size=10, color=ACCENT_R if ci == 0 else LGRAY, bold=(ci == 0))
        else:
            cell_fill(c, bg)
            cell_text(c, val, size=10, bold=(ci == 0), color=WHITE if ci == 0 else LGRAY)

label(sl, "Override rule: Any MCP server connecting to a live system auto-escalates to Zone 3, regardless of what data you think you're working with.", 0.3, 6.48, 12.73, size=10, color=MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LICENSING DECISION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Licensing & Tooling")
slide_title(sl, "Licensing Decision")
divider_line(sl)

# Three columns: NOW / GITHUB TEAMS / FUTURE
cols_data = [
    {
        "title": "NOW",
        "sub": "Claude Teams Plan",
        "color": BLUE,
        "items": [
            "$30 / seat / month",
            "Already contracted — move immediately",
            "Sandbox (Seatbelt / bubblewrap)",
            "managed-settings.json controls",
            "managed-mcp.json allowlist",
            "SSO / SAML via Azure AD or Okta",
            "Usage analytics (per-user, CSV export)",
            "Hooks for audit logging",
        ],
    },
    {
        "title": "ALSO NOW",
        "sub": "GitHub Teams Plan",
        "color": RGBColor(0x1C, 0x7A, 0x3E),
        "items": [
            "$4 / seat / month",
            "Version control for all citizen work",
            "180-day audit log (org-level)",
            "Branch protection + PR reviews",
            "Highlander webhook ingestion point",
            "Compensating SOC2 audit trail",
            "Rollback for any AI mistake",
        ],
    },
    {
        "title": "FUTURE (After Q1)",
        "sub": "Anthropic Enterprise API",
        "color": GOLD,
        "items": [
            "Per-user token-level usage tracking",
            "Native audit logs (180-day export)",
            "Compliance API → SIEM integration",
            "SCIM auto-provisioning/deprovisioning",
            "Custom data retention (ZDR option)",
            "Advanced admin controls + RBAC",
            "Evaluate once usage patterns are known",
        ],
    },
]

cw = 3.9
cx = 0.35
for i, col in enumerate(cols_data):
    bx = cx + i * (cw + 0.22)
    rect(sl, bx, 1.5, cw, 0.42, col["color"])
    tb = box(sl, bx, 1.5, cw, 0.42)
    tf = tb.text_frame
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, col["title"], size=14, bold=True, color=WHITE)

    rect(sl, bx, 1.92, cw, 0.32, ROW_ALT)
    tb2 = box(sl, bx, 1.92, cw, 0.32)
    tf2 = tb2.text_frame
    tf2.margin_top = Inches(0.05)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    add_run(p2, col["sub"], size=12, color=col["color"], italic=True)

    rect(sl, bx, 2.24, cw, 4.2, ROW_EVEN)
    tb3 = box(sl, bx + 0.1, 2.32, cw - 0.2, 4.0)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    first = True
    for item in col["items"]:
        p3 = tf3.paragraphs[0] if first else para_in(tf3)
        first = False
        p3.space_after = Pt(5)
        add_run(p3, "• " + item, size=11, color=LGRAY)

label(sl, "Why Teams is sufficient to start: sandbox + managed settings + hooks + SSO are all available on the Teams plan. GitHub fills the audit trail gap.", 0.35, 6.65, 12.63, size=10, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLAUDE DESKTOP vs CLAUDE CODE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Licensing & Tooling")
slide_title(sl, "Claude Desktop vs Claude Code")
divider_line(sl)

label(sl, "Different tools for different users. Both connect to MCP servers. Only Claude Code has the full managed settings governance stack.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

# Two panels
for i, (ptitle, pcolor, rows) in enumerate([
    ("Claude Desktop", BLUE, [
        ("Who", "Non-technical citizens — no terminal required"),
        ("Interface", "GUI chat app with visual diffs, drag-and-drop"),
        ("Setup", "Download → sign in → select folder → describe task"),
        ("MCP", "One-click Desktop Extensions (.mcpb files)"),
        ("Workflow", "Conversational, iterative, visual approval of changes"),
        ("Background work", "Cowork tab — spawn agents that run while you work"),
        ("Governance", "Limited managed settings; relies on GitHub as guardrail"),
        ("Best for", "Reports, data pipelines, automations, business analysis"),
    ]),
    ("Claude Code (Terminal)", RGBColor(0x1A, 0x48, 0xBB), [
        ("Who", "More technical citizens comfortable with the terminal"),
        ("Interface", "CLI — runs in terminal, diffs shown in text"),
        ("Setup", "Terminal → navigate to project folder → run `claude`"),
        ("MCP", "Via settings.json (JSON config, more control)"),
        ("Workflow", "Full Git automation — branch, commit, PR without GitHub UI"),
        ("Background work", "Full shell automation, scripts, CI/CD integration"),
        ("Governance", "Full managed-settings.json stack — strongest compliance posture"),
        ("Best for", "Complex projects, multi-file work, automated pipelines"),
    ]),
]):
    bx = 0.3 + i * 6.52
    rect(sl, bx, 1.6, 6.2, 0.38, pcolor)
    tb = box(sl, bx, 1.6, 6.2, 0.38)
    tf = tb.text_frame
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, ptitle, size=15, bold=True, color=WHITE)

    for ri, (label_txt, val_txt) in enumerate(rows):
        ry = 2.05 + ri * 0.54
        bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
        rect(sl, bx, ry, 6.2, 0.52, bg)
        label(sl, label_txt, bx + 0.1, ry + 0.08, 1.35, size=10, bold=True, color=pcolor)
        label(sl, val_txt, bx + 1.5, ry + 0.08, 4.55, size=10, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GITHUB TEAMS vs NO GITHUB
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Licensing & Tooling")
slide_title(sl, "GitHub Teams vs No GitHub")
divider_line(sl)

label(sl, "GitHub Teams IS the compliance infrastructure on the Teams plan. Without it, we have no audit trail, no rollback, no Highlander tracking.", 0.3, 1.35, 12.73, size=12, color=BLUE, bold=True)

headers_row = ["Capability", "With GitHub Teams  ($4/seat/mo)", "Without GitHub"]
tbl = make_table(sl, 10, 3, 0.3, 1.72, 12.73, 4.95, col_widths=[3.0, 4.87, 4.86])
header_row(tbl, headers_row)

rows_data = [
    ("Version control & rollback", "Every change tracked. Rollback any AI mistake instantly.", "No rollback. AI errors are permanent."),
    ("Audit trail", "GitHub org audit log — 180-day retention. Owner-accessible.", "No audit trail. Cannot demonstrate controls to auditors."),
    ("Branch protection", "Require PR reviews before merge. Block direct pushes to main.", "No protection. Anyone can push anything directly."),
    ("Code review gate", "Human approval required before AI-generated work is merged.", "No review gate. AI output goes directly into use."),
    ("Highlander tracking", "GitHub webhook → Highlander ODL → commit velocity, AI co-auth, OKR attribution.", "Cannot connect to Highlander. No business value measurement."),
    ("SOC2 compensating control", "GitHub log + hook logging = defensible CC7.2 compensating control.", "No path to SOC2 compliance on Teams plan."),
    ("Team collaboration", "PRs, issues, project boards. Non-technical reviewers can participate.", "Siloed. Work stays on individual machines."),
    ("CI/CD & automation", "GitHub Actions can run tests, lint, deploy on every PR.", "No automation. Manual deployment only."),
    ("Cost", "$100/month for 25 seats.", "$0 — but no compliance story, no tracking, no rollback."),
]
for ri, row in enumerate(rows_data):
    bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        cell_fill(c, bg)
        clr = ACCENT_R if (ci == 2 and ri <= 5) else (LGRAY if ci != 0 else WHITE)
        cell_text(c, val, size=9.5, bold=(ci == 0), color=clr)

label(sl, "Bottom line: $4/seat/month is not optional infrastructure — it is the SOC2 compliance layer for the Teams plan.", 0.3, 6.78, 12.73, size=10, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TEAMS PLAN CONTROL INVENTORY
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Technical Guardrails")
slide_title(sl, "What We Control on Claude Teams Plan")
divider_line(sl)

controls = [
    ("managed-settings.json", "Permission deny rules, command restrictions, file access limits, network domain allowlist. Deployed by IT — users cannot override.", "Prevents Claude from accessing sensitive paths or running destructive commands."),
    ("managed-mcp.json", "Empty allowlist = lockdown mode. Only vetted MCP servers permitted. Deployed to system directories. Users cannot add unapproved servers.", "Eliminates MCP supply-chain risk (tool-poisoning attacks, rug pulls, cross-server exfiltration)."),
    ("Global CLAUDE.md", "Governance rules injected into every session. Jira-first enforcement, no-credential rules, Driving Zone reminders. Users see it but cannot remove it.", "Soft guardrail in every session. Primes the model before any task begins."),
    ("Hooks — PreToolUse", "Block: rm -rf, push to main/master, credential/secret patterns in commands. Shell commands evaluated before execution.", "Hard stop on the most common accidental or malicious destructive actions."),
    ("Hooks — PostToolUse", "Audit log: every tool use timestamped and written to log file. HTTP hook option: POST JSON to central endpoint for SIEM ingestion.", "Creates the compensating audit trail required for SOC2 CC7.2 on the Teams plan."),
    ("Sandbox (Seatbelt/bwrap)", "OS-level filesystem isolation (write only to allowed dirs) + network proxy (approved domains only). Cannot be bypassed by Claude prompts.", "The only true enforcement layer. Even a compromised session cannot escape the sandbox."),
    ("SSO / SAML", "Available on Teams. Integrates with Azure AD, Okta, Google Workspace. JIT provisioning on first login. Auth logging via your IdP.", "Centralises access control. IdP logs provide authentication audit trail even without Claude native audit logs."),
    ("Usage Analytics", "Per-user, per-model, per-day metrics. Suggestion acceptance rate, lines accepted. CSV export. ~1-day lag.", "Not audit logs, but a program health signal. Shows who is active, which models are used."),
]

ty = 1.58
for i, (ctrl, desc, why) in enumerate(controls):
    bg = ROW_ALT if i % 2 == 0 else ROW_EVEN
    rect(sl, 0.3, ty, 12.73, 0.6, bg)
    label(sl, ctrl, 0.42, ty + 0.08, 2.3, size=10, bold=True, color=BLUE)
    label(sl, desc, 2.78, ty + 0.05, 5.55, size=9, color=LGRAY)
    label(sl, why, 8.42, ty + 0.05, 4.5, size=9, color=MGRAY, italic=True)
    ty += 0.63

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GAPS & COMPENSATING CONTROLS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Technical Guardrails")
slide_title(sl, "What We Don't Have — and How We Compensate")
divider_line(sl)

label(sl, "Teams plan lacks Enterprise-only detective controls. Each gap has a documented compensating control. Auditors accept compensating controls when they are intentional, documented, and consistently applied.", 0.3, 1.35, 12.73, size=11, color=MGRAY, italic=True)

tbl = make_table(sl, 6, 3, 0.3, 1.78, 12.73, 4.3, col_widths=[2.8, 4.5, 5.43])
header_row(tbl, ["Enterprise-Only Feature", "Our Compensating Control", "SOC2 Criteria Addressed"])

gaps = [
    ("Native audit logs\n(180-day, JSON/CSV export)", "GitHub org audit log (180-day) + PostToolUse hook logging to local file or HTTP endpoint. Stream GitHub logs to S3/Datadog for 12-month retention.", "CC7.2 — System monitoring.\nCC6.1 — Logical access."),
    ("SCIM auto-deprovisioning", "SSO JIT provisioning + manual offboarding checklist (IT removes GitHub + Claude seat on HR termination event).", "CC6.1 — Logical access provisioning.\nCC6.2 — Timely deprovisioning."),
    ("Compliance API\n(programmatic access)", "Highlander GitHub ingestion provides commit-level activity data. Usage analytics CSV exported weekly for programme records.", "CC7.2 — Monitoring.\nPI1.4 — Processing integrity."),
    ("Custom data retention / ZDR", "Documented 90-day Teams retention policy + GitHub repo as permanent record. Accepted as proportional for our risk profile.", "ISO 27001 A.8.10 — Data lifecycle."),
    ("IP allowlisting", "Network sandbox (managed-settings.json) restricts Claude's outbound access. VPN enforced for Zone 3 work.", "C1.1 — Confidentiality."),
]
for ri, row in enumerate(gaps):
    bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        cell_fill(c, bg)
        cell_text(c, val, size=9.5, color=LGRAY if ci != 0 else WHITE, bold=(ci == 0))

rect(sl, 0.3, 6.25, 12.73, 0.55, RGBColor(0x0A, 0x25, 0x5A))
tb = box(sl, 0.42, 6.28, 12.5, 0.5)
tf = tb.text_frame
p = tf.paragraphs[0]
add_run(p, "Good faith framing: ", size=11, bold=True, color=BLUE)
add_run(p, "We have documented every gap, its compensating control, and its SOC2 mapping. Auditors evaluate proportionality + consistency — not perfection. This is sufficient.", size=11, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SOC2 MINIMUM vs MAXIMUM
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "SOC2 Compliance")
slide_title(sl, "SOC2 Compliance — Minimum vs Maximum")
divider_line(sl)

label(sl, "SOC2 auditors evaluate proportionality and consistency — not whether you have every possible control. Good faith + documentation = defensible.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

# Two panels
for i, (ptitle, pcolor, badge, items) in enumerate([
    ("MINIMUM — Good Faith (Where We Are)", BLUE, "OUR CURRENT STATE", [
        "AI Acceptable Use Policy — documented, signed by all employees",
        "Training records — completion timestamps in Highlander Learning Platform",
        "GitHub audit logs — 180-day retention, org-level, owner-accessible",
        "Hook-based logging — every Claude tool-use timestamped locally",
        "Signed AUP acknowledgments — per employee, dated, on file",
        "Branch protection — human review required before any AI code merges (CC8.1)",
        "Driving Zones — documented RBAC for data access tiers (CC6.3)",
        "SSO via IdP — centralised authentication logging (CC6.1)",
        "Sandbox — OS-level enforcement preventing data exfiltration (C1.1)",
    ]),
    ("MAXIMUM — Full Detective Controls (Future State)", GOLD, "ROADMAP TARGET", [
        "Claude Enterprise API — native audit logs, 180-day CSV export",
        "Compliance API → SIEM (Splunk / Datadog / Elastic) — real-time monitoring",
        "SCIM provisioning — automated deprovisioning on HR termination event",
        "Hash-chained immutable audit logs — tamper-evident log integrity",
        "Real-time anomaly alerting — unusual usage triggers security review",
        "Zero Data Retention (ZDR) — Claude never persists conversation data",
        "Custom data retention policies — ISO 27001 A.8.10 fully addressed",
        "IP allowlisting — restricts Claude access to known corporate networks",
        "HIPAA-ready configuration — required only if health data in scope",
    ]),
]):
    bx = 0.3 + i * 6.52
    rect(sl, bx, 1.62, 6.2, 0.36, pcolor)
    tb = box(sl, bx, 1.62, 6.2, 0.36)
    tf = tb.text_frame
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, ptitle, size=11, bold=True, color=WHITE)

    for ri, item in enumerate(items):
        ry = 2.04 + ri * 0.52
        bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
        rect(sl, bx, ry, 6.2, 0.5, bg)
        check = "✓" if i == 0 else "→"
        clr = ACCENT_G if i == 0 else GOLD
        label(sl, check, bx + 0.08, ry + 0.1, 0.28, size=11, bold=True, color=clr)
        label(sl, item, bx + 0.4, ry + 0.08, 5.65, size=9.5, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SOC2 CONTROL MAP
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "SOC2 Compliance")
slide_title(sl, "Our SOC2 Control Map")
divider_line(sl)

tbl = make_table(sl, 7, 3, 0.3, 1.55, 12.73, 5.2, col_widths=[1.9, 4.6, 6.23])
header_row(tbl, ["SOC2 Criteria", "Control Required", "How We Address It (Teams Plan)"])

controls_map = [
    ("CC6.1\nLogical Access", "Document who can use Claude, approval process, monitoring.", "SSO (Azure AD / Okta) + JIT provisioning. Certification gate before seat issued. Manual offboarding checklist."),
    ("CC6.3\nRole-Based Access", "RBAC for AI tool access. Least-privilege for code repos.", "Driving Zones = tiered data access model. Zone 3 requires manager approval. managed-mcp.json = deny-by-default."),
    ("CC7.2\nSystem Monitoring", "Audit logs of all AI tool usage. SIEM integration.", "GitHub audit log (180-day) + hook-based PostToolUse logging. Highlander commit metrics. Usage analytics CSV."),
    ("CC8.1\nChange Management", "Code review for AI-generated code. Human validation.", "Branch protection on all citizen repos. PR required before merge. Zone 2 minimum process enforces this."),
    ("PI1.4\nProcessing Integrity", "Demonstrate AI outputs reviewed for accuracy.", "Human review gate at Zone 2+. Zone 3 requires manager sign-off. No autonomous production deploys without review."),
    ("C1.1\nConfidentiality", "Classify data before sending to AI. Encryption. DLP.", "Sandbox network isolation (approved domains only). MCP allowlist. AUP data classification rules. No Zone 3 without Jira ticket."),
]
for ri, row in enumerate(controls_map):
    bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        cell_fill(c, bg)
        cell_text(c, val, size=9.5, bold=(ci == 0), color=BLUE if ci == 0 else LGRAY)

label(sl, "Note: ISO 27001 A.5.10, A.8.15, A.8.16 are addressed by the same controls above. The mapping is documented in the compliance workbook.", 0.3, 6.88, 12.73, size=9, color=MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PROJECT LIFECYCLE (BOTH PATHS)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Workflow Walkthroughs")
slide_title(sl, "Project Lifecycle — The Two Paths")
divider_line(sl)

# Zone 2 path
label(sl, "ZONE 2 PATH — Internal, non-confidential data", 0.3, 1.42, 12.73, size=12, bold=True, color=BLUE)
z2_steps = ["Idea", "GitHub\nbranch", "Claude\nbuilds it", "Review\n& iterate", "PR +\ncode review", "Merge", "Highlander\ningests", "OKR\nattributed"]
z2_bw, z2_bh, z2_by = 1.45, 0.72, 1.75
z2_x = 0.3
for i, step in enumerate(z2_steps):
    bx = z2_x + i * (z2_bw + 0.1)
    rect(sl, bx, z2_by, z2_bw, z2_bh, BLUE if i < 6 else RGBColor(0x1A, 0x48, 0xBB))
    tb = box(sl, bx, z2_by, z2_bw, z2_bh)
    tf = tb.text_frame
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, step, size=11, bold=True, color=WHITE)
    if i < len(z2_steps) - 1:
        label(sl, "→", bx + z2_bw + 0.01, z2_by + 0.22, 0.1, size=13, color=BLUE, bold=True)

label(sl, "No Jira ticket required. GitHub branch + code review is sufficient. Claude Desktop or Claude Code both work.", 0.3, 2.58, 12.73, size=10, color=MGRAY, italic=True)

# Zone 3 path
label(sl, "ZONE 3 PATH — Confidential / PII / client data / any MCP server", 0.3, 3.3, 12.73, size=12, bold=True, color=GOLD)
z3_steps = ["Jira\nticket", "Manager\napproval", "Approved\nMCP only", "Claude\nbuilds it", "Zone 3\naudit log", "PR +\nmanager review", "Merge", "Highlander\ningests"]
z3_by = 3.62
for i, step in enumerate(z3_steps):
    bx = z2_x + i * (z2_bw + 0.1)
    rect(sl, bx, z3_by, z2_bw, z2_bh, RGBColor(0x5A, 0x3A, 0x00) if i < 3 else (BLUE if i < 7 else RGBColor(0x1A, 0x48, 0xBB)))
    tb = box(sl, bx, z3_by, z2_bw, z2_bh)
    tf = tb.text_frame
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, step, size=11, bold=True, color=GOLD if i < 3 else WHITE)
    if i < len(z3_steps) - 1:
        label(sl, "→", bx + z2_bw + 0.01, z3_by + 0.22, 0.1, size=13, color=GOLD, bold=True)

label(sl, "Zone 3 is triggered automatically by any MCP server connecting to a live system — regardless of what data the citizen thinks they're working with.", 0.3, 4.44, 12.73, size=10, color=MGRAY, italic=True)

# METADATA.yaml callout
rect(sl, 0.3, 5.05, 12.73, 1.1, ROW_ALT)
tb_m = box(sl, 0.45, 5.1, 12.43, 1.0)
tf_m = tb_m.text_frame
tf_m.word_wrap = True
p_m = tf_m.paragraphs[0]
add_run(p_m, "Every citizen repo includes METADATA.yaml: ", size=11, bold=True, color=BLUE)
add_run(p_m, "owner, team, OKR reference, business problem, and zone. This is what links a GitHub commit to an OKR in Highlander — no manual reporting.", size=11, color=LGRAY)
p_m2 = para_in(tf_m)
p_m2.space_before = Pt(4)
add_run(p_m2, "Repo naming: citizen-<team>-<project>   |   GitHub topic tag: citizen-ai   |   Company email required for identity mapping in Highlander", size=10, color=MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CITIZEN WORKFLOW: CLAUDE DESKTOP
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Workflow Walkthroughs")
slide_title(sl, "Citizen Workflow: Claude Desktop (No Terminal Required)")
divider_line(sl)

label(sl, "Recommended path for non-technical citizens. GUI-driven, visual, conversational. No Git or terminal knowledge needed.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

steps_d = [
    ("1", "Open Claude Desktop", "Download the app, sign in with company SSO, select your project folder (e.g., ~/citizen-ops-deal-tracker)."),
    ("2", "Describe your task", "Type in plain English: \"Build a weekly ARR report from our Google Sheet — pull data, calculate YTD, email to sales@helmmarkets.com.\""),
    ("3", "Claude builds it", "Claude reads your folder, writes the code, and shows visual diffs — exactly what changed, highlighted in color. You see every line before it's saved."),
    ("4", "Review & iterate", "Approve each change or say \"make this metric red if it drops more than 10%.\" Completely conversational. No coding required."),
    ("5", "Commit to GitHub branch", "Click the GitHub button or ask Claude: \"Commit this to a new branch called citizen-ops-arr-report.\" Branch created, code pushed."),
    ("6", "PR → Review → Merge → Tracked", "PR opens automatically. Teammate reviews. On merge, Highlander ingests the commit. OKR attributed via METADATA.yaml."),
]

sw, sh = 12.73, 0.78
sx = 0.3
for i, (num, title, body) in enumerate(steps_d):
    sy = 1.72 + i * (sh + 0.05)
    bg = ROW_ALT if i % 2 == 0 else ROW_EVEN
    rect(sl, sx, sy, sw, sh, bg)
    rect(sl, sx, sy, 0.42, sh, BLUE)
    tb_n = box(sl, sx + 0.05, sy + 0.18, 0.32, 0.4)
    p_n = tb_n.text_frame.paragraphs[0]
    p_n.alignment = PP_ALIGN.CENTER
    add_run(p_n, num, size=14, bold=True, color=WHITE)
    label(sl, title, sx + 0.55, sy + 0.1, 2.5, size=12, bold=True, color=WHITE)
    label(sl, body, sx + 3.15, sy + 0.1, 9.75, size=10.5, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CITIZEN WORKFLOW: CLAUDE CODE IN TERMINAL
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Workflow Walkthroughs")
slide_title(sl, "Citizen Workflow: Claude Code in Terminal")
divider_line(sl)

label(sl, "For citizens with some technical comfort. Full Git automation. Stronger compliance posture via managed-settings.json.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

steps_c = [
    ("1", "Open terminal & navigate", "cd ~/citizen-finance-arr-report   then run:   claude   (the CLI launches automatically)"),
    ("2", "Describe your task", "Type in natural language — same as Desktop. Claude reads your entire codebase context before responding."),
    ("3", "Claude proposes changes", "Claude shows diffs in the terminal. Every file change is explicit. You approve or reject each one with a keystroke."),
    ("4", "Review diffs", "Approve changes interactively. Claude never modifies files you haven't approved. The sandbox prevents writes outside your project folder."),
    ("5", "Claude handles Git automatically", "On approval: Claude creates a branch (citizen-finance-arr-report), commits with a descriptive message, and opens a PR — no GitHub UI needed."),
    ("6", "PR → Review → Merge → Tracked", "Same as Desktop: PR reviewed by teammate, merged, Highlander ingests, OKR attributed via METADATA.yaml."),
]

for i, (num, title, body) in enumerate(steps_c):
    sy = 1.72 + i * 0.83
    bg = ROW_ALT if i % 2 == 0 else ROW_EVEN
    rect(sl, 0.3, sy, 12.73, 0.78, bg)
    rect(sl, 0.3, sy, 0.42, 0.78, RGBColor(0x1A, 0x48, 0xBB))
    tb_n = box(sl, 0.35, sy + 0.18, 0.32, 0.4)
    p_n = tb_n.text_frame.paragraphs[0]
    p_n.alignment = PP_ALIGN.CENTER
    add_run(p_n, num, size=14, bold=True, color=WHITE)
    label(sl, title, 0.85, sy + 0.1, 2.5, size=12, bold=True, color=WHITE)
    label(sl, body, 3.45, sy + 0.1, 9.45, size=10.5, color=LGRAY)

label(sl, "Requires terminal comfort. Recommend only for citizens who've completed Module 16 (GitHub) and are comfortable navigating file paths. Not required — Desktop is fully sufficient for most citizen work.", 0.3, 6.78, 12.73, size=9.5, color=MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — HIGHLANDER INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Highlander Integration")
slide_title(sl, "Highlander — How We Prove Business Value")
divider_line(sl)

# Flow diagram
flow = ["Citizen\nCommits", "GitHub\nWebhook", "Highlander\nODL", "Analyzers\nRun", "Metrics\nDashboard", "OKR\nProof"]
fw, fh, fy = 1.85, 0.8, 1.58
ftotal = len(flow) * fw + (len(flow) - 1) * 0.35
fx_start = (13.33 - ftotal) / 2
for i, step in enumerate(flow):
    bx = fx_start + i * (fw + 0.35)
    clr = BLUE if i % 2 == 0 else RGBColor(0x1A, 0x48, 0xBB)
    rect(sl, bx, fy, fw, fh, clr)
    tb = box(sl, bx, fy, fw, fh)
    tf = tb.text_frame
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, step, size=12, bold=True, color=WHITE)
    if i < len(flow) - 1:
        label(sl, "→", bx + fw + 0.04, fy + 0.24, 0.28, size=16, color=BLUE, bold=True)

# Metrics tracked
label(sl, "What Highlander tracks from citizen activity:", 0.3, 2.6, 12.73, size=12, bold=True, color=WHITE)

metrics = [
    ("Commit velocity\nper citizen", "Who is active vs. inactive. Adoption rate over time."),
    ("AI co-authorship %", "Force multiplier — % of commits where Claude was the author."),
    ("Gap ratio reduction", "Are citizen tools reducing operational bottlenecks?"),
    ("Constraint reduction", "Is support case volume falling where citizens have built tools?"),
    ("OKR attribution", "METADATA.yaml links every commit to a specific OKR."),
    ("Cert completion\nfunnel", "Enrolled → Certified → First commit → Active → Driving value."),
]

mw = 2.0
mx = 0.3
for i, (metric, desc) in enumerate(metrics):
    bx = mx + i * (mw + 0.12)
    rect(sl, bx, 2.98, mw, 1.5, ROW_ALT)
    rect(sl, bx, 2.98, mw, 0.06, BLUE)
    label(sl, metric, bx + 0.08, 3.06, mw - 0.16, size=10, bold=True, color=BLUE)
    label(sl, desc, bx + 0.08, 3.42, mw - 0.16, size=9.5, color=LGRAY)

# Setup callout
rect(sl, 0.3, 4.65, 12.73, 0.78, RGBColor(0x0A, 0x25, 0x5A))
tb2 = box(sl, 0.45, 4.7, 12.43, 0.68)
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
add_run(p2, "One-time setup: ~8 hours.  ", size=12, bold=True, color=BLUE)
add_run(p2, "Register repos in admin.github_repositories → map identities in engineer_email_mapping → add 'citizen' filter to engineer_effectiveness analyzer → deploy repo template.", size=11, color=LGRAY)

# Board narrative
rect(sl, 0.3, 5.55, 12.73, 0.78, ROW_ALT)
tb3 = box(sl, 0.45, 5.6, 12.43, 0.68)
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
add_run(p3, "Board narrative: ", size=12, bold=True, color=GOLD)
add_run(p3, "\"X citizen developers shipped Y tools → closed Z operational bottlenecks → reduced support case volume by N% against OKR [X].\" Highlander generates this automatically.", size=11, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — COURSE & CERTIFICATION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Course & Certification")
slide_title(sl, "Certification — The Driver's License Gate")
divider_line(sl)

label(sl, "No license without the course. The Highlander Learning Platform (already built) tracks every completion in quiz_results.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

# Funnel
funnel_steps = [
    ("Enrolled", MGRAY, "Seat assigned, course link sent"),
    ("Certified", BLUE, "All 16 modules + AUP signed"),
    ("First Commit", RGBColor(0x1A, 0x7A, 0x3E), "First citizen repo created"),
    ("Active", ACCENT_G, "Commits every 2 weeks"),
    ("Driving Value", GOLD, "Highlander shows OKR attribution"),
]
fw2 = 12.73
for i, (step, clr, desc) in enumerate(funnel_steps):
    fi_w = fw2 - i * 1.1
    fi_x = 0.3 + i * 0.55
    fy2  = 1.72 + i * 0.72
    rect(sl, fi_x, fy2, fi_w, 0.6, RGBColor(max(clr[0] - 60, 0), max(clr[1] - 60, 0), max(clr[2] - 60, 0)))
    rect(sl, fi_x, fy2, fi_w, 0.04, clr)
    label(sl, step, fi_x + 0.15, fy2 + 0.1, 2.8, size=12, bold=True, color=clr)
    label(sl, desc, fi_x + 3.1, fy2 + 0.1, fi_w - 3.3, size=11, color=LGRAY)

# 16 modules grid
label(sl, "16-module curriculum in Highlander Learning Platform:", 0.3, 5.42, 12.73, size=11, bold=True, color=WHITE)

modules_short = [
    "01 What Claude Is", "02 How Models Work", "03 Tokens", "04 Context Window",
    "05 Specificity", "06 Anatomy of a Prompt", "07 Iteration", "08 Tips & Tricks",
    "09 Operating Framework", "10 Projects", "11 Docs & Multimodal", "12 Team Use Cases",
    "13 Safety & Responsible Use", "14 MCP & Agents", "15 GitHub", "16 Databases",
]
mw2, mh2 = 2.95, 0.32
for i, mod in enumerate(modules_short):
    col = i % 4
    row = i // 4
    bx2 = 0.3 + col * (mw2 + 0.13)
    by2 = 5.78 + row * (mh2 + 0.05)
    bg = BLUE if mod in ["09 Operating Framework", "13 Safety & Responsible Use", "14 MCP & Agents", "15 GitHub"] else ROW_ALT
    rect(sl, bx2, by2, mw2, mh2, bg)
    label(sl, mod, bx2 + 0.1, by2 + 0.06, mw2 - 0.2, size=9, color=WHITE if bg == BLUE else LGRAY, bold=(bg == BLUE))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — MEETING CADENCE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Program Cadence")
slide_title(sl, "Meeting Cadence — Culture + Accountability")
divider_line(sl)

label(sl, "Accountability without bureaucracy. Three recurring touchpoints keep the program healthy without overwhelming citizens.", 0.3, 1.35, 12.73, size=12, color=MGRAY, italic=True)

meetings = [
    ("Weekly", "Open Office Hours", "Optional · 30 min", BLUE, [
        "Drop-in Q&A — no agenda required",
        "Unblock citizens on tasks or zone questions",
        "Share wins and showcase cool outputs",
        "Pairs with async Slack channel for between-session questions",
        "No required attendance — but keeps the door open",
    ]),
    ("Monthly", "All-Hands Program Call", "Mandatory · 30 min", RGBColor(0x1A, 0x48, 0xBB), [
        "Highlander metrics recap — who's active, OKRs being driven",
        "One citizen showcase: real tool built in the last month",
        "Policy/tooling updates — newly approved MCP servers, zone clarifications",
        "10 min open Q&A",
        "Recorded for anyone who misses",
    ]),
    ("Quarterly", "1:1 Check-ins", "Mandatory · 20 min each", GOLD, [
        "OKR review: Highlander data for that citizen's repos",
        "Zone compliance check: any edge cases or near-misses?",
        "Coaching: better workflows, zone upgrades if ready",
        "Flag policy gaps → feed into curriculum updates",
        "Re-certification reminder ahead of annual review",
    ]),
]

cw2 = 3.95
for i, (freq, title, sub, clr, items) in enumerate(meetings):
    bx2 = 0.3 + i * (cw2 + 0.27)
    rect(sl, bx2, 1.65, cw2, 0.38, clr)
    tb = box(sl, bx2, 1.65, cw2, 0.38)
    tf = tb.text_frame
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, f"{freq}  —  {title}", size=13, bold=True, color=WHITE)

    rect(sl, bx2, 2.03, cw2, 0.28, ROW_ALT)
    tb2 = box(sl, bx2, 2.03, cw2, 0.28)
    tf2 = tb2.text_frame
    tf2.margin_top = Inches(0.04)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    add_run(p2, sub, size=10, color=clr, italic=True)

    rect(sl, bx2, 2.31, cw2, 3.65, ROW_EVEN)
    tb3 = box(sl, bx2 + 0.1, 2.42, cw2 - 0.2, 3.45)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    first = True
    for item in items:
        p3 = tf3.paragraphs[0] if first else para_in(tf3)
        first = False
        p3.space_after = Pt(7)
        add_run(p3, "• " + item, size=11, color=LGRAY)

# Ad hoc
rect(sl, 0.3, 6.1, 12.73, 0.6, RGBColor(0x2A, 0x08, 0x08))
tb4 = box(sl, 0.45, 6.13, 12.43, 0.5)
tf4 = tb4.text_frame
p4 = tf4.paragraphs[0]
add_run(p4, "Ad Hoc — Incident Review: ", size=11, bold=True, color=ACCENT_R)
add_run(p4, "Triggered by any No-Drive Zone violation. Program lead + manager + IT review within 48 hours. Outcome: re-certification, remediation, or license revocation.", size=11, color=LGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — COST MODEL & ROI
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
section_pill(sl, "Cost & ROI")
slide_title(sl, "Cost Model & Return on Investment (25 Citizens)")
divider_line(sl)

# Cost table
tbl = make_table(sl, 7, 4, 0.3, 1.55, 7.8, 3.85, col_widths=[3.2, 1.6, 1.5, 1.5])
header_row(tbl, ["Category", "Monthly", "Annual", "Status"])

cost_rows = [
    ("Claude Teams (25 seats @ $30/seat)", "$750", "$9,000", "Already contracted"),
    ("GitHub Teams (25 seats @ $4/seat)", "$100", "$1,200", "New spend"),
    ("Highlander infra (marginal increase)", "$8", "$96", "Already running"),
    ("Program management (~5 hrs/month)", "$750", "$9,000", "Internal cost"),
    ("One-time setup (engineering + course)", "—", "—", "$6,150 one-time"),
    ("TOTAL — Year 1", "$1,608", "$19,296", "~$25K all-in"),
]
for ri, row in enumerate(cost_rows):
    bg = ROW_ALT if ri % 2 == 0 else ROW_EVEN
    is_total = (ri == 5)
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        cell_fill(c, BLUE if is_total else bg)
        cell_text(c, val, size=10.5, bold=is_total, color=WHITE if is_total else (LGRAY if ci != 0 else WHITE))

# ROI panel
rect(sl, 8.3, 1.55, 4.73, 3.85, ROW_ALT)
rect(sl, 8.3, 1.55, 4.73, 0.06, BLUE)
tb_roi = box(sl, 8.42, 1.65, 4.5, 3.65)
tf_roi = tb_roi.text_frame
tf_roi.word_wrap = True

roi_items = [
    ("ROI MODEL", None, BLUE, 13),
    ("25 people × 1 hr/week saved", None, LGRAY, 11),
    ("× 50 weeks × $75/hr", None, LGRAY, 11),
    ("= $93,750 recaptured capacity", None, WHITE, 12),
    ("", None, LGRAY, 8),
    ("Program cost (year 1): ~$25K", None, MGRAY, 10),
    ("Break-even: ~3 months", None, ACCENT_G, 12),
    ("Year 1 ROI: 2.8×", None, ACCENT_G, 14),
    ("", None, LGRAY, 8),
    ("+ Highlander data proves", None, MGRAY, 10),
    ("constraint reduction & OKR", None, MGRAY, 10),
    ("impact independently.", None, MGRAY, 10),
]
first = True
for text, _, clr, sz in roi_items:
    p = tf_roi.paragraphs[0] if first else para_in(tf_roi)
    first = False
    p.space_after = Pt(2)
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=sz, color=clr, bold=(sz >= 13))

# Future state
rect(sl, 0.3, 5.55, 12.73, 0.78, RGBColor(0x2A, 0x1A, 0x00))
tb5 = box(sl, 0.45, 5.6, 12.43, 0.68)
tf5 = tb5.text_frame
tf5.word_wrap = True
p5 = tf5.paragraphs[0]
add_run(p5, "Future — Anthropic Enterprise API: ", size=11, bold=True, color=GOLD)
add_run(p5, "adds ~$1,125–1,875/month incremental (delta from Teams to Enterprise pricing). Evaluate after Q1 once usage patterns and SOC2 audit findings are known.", size=11, color=LGRAY)

# ─── Save ────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✓  Deck saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
